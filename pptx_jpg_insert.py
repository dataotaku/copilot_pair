from pptx import Presentation
from pptx.util import Inches
import os
from PIL import Image

# # Presentation 객체 생성
prs = Presentation()

# 사용자 정의 용지 크기 설정 (A4 용지 가로 11.69인치, 세로 8.27인치)


# for lo in prs.slide_layouts:
#     print(lo.name)

# Blank 슬라이드 레이아웃 가져오기
blank_slide_layout = prs.slide_layouts[6]

# 슬라이드 추가
slide = prs.slides.add_slide(blank_slide_layout)

# 제목과 부제목 설정
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

# title.text = "Hello, python-pptx!"
# subtitle.text = "python-pptx 라이브러리를 사용하여 PowerPoint 파일을 생성합니다."

# 새 슬라이드 추가 (빈 슬라이드)
# blank_slide_layout = prs.slide_layouts[6]
# slide = prs.slides.add_slide(blank_slide_layout)

# 이미지 추가
# img_path = "path/to/image.jpg"
# left = Inches(2)
# top = Inches(2)
# slide.shapes.add_picture(img_path, left, top)

# 이미지가 있는 폴더 지정
img_folder = "c:/Users/unhoc/Documents/copilot_pair/AI코딩_테스트"
path_dir = "c:/Users/unhoc/Documents/copilot_pair/pptx"
temp_img_path = "temp_image.jpg"
prs = Presentation()
# print(type(os.listdir(img_folder)))
a4_width_inch = 8.27  # 8.27
a4_height_inch = 11.69  # 11.69
# 폴더 내의 모든 파일 순회
prs.slide_width = Inches(a4_width_inch)
prs.slide_height = Inches(a4_height_inch)
img_files = os.listdir(img_folder)
img_files.reverse()

for img_name in img_files:
    img_path = os.path.join(img_folder, img_name)
    if os.path.isfile(img_path):
        # 이미지 열기
        picture_path = str(os.path.join(img_folder, img_name)).replace("\\", "/")
        # 새 슬라이드 추가 (빈 슬라이드)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # 이미지 추가
        left = Inches(0)
        top = Inches(0)
        slide.shapes.add_picture(picture_path, left, top)

        # 새 슬라이드 추가 (빈 슬라이드)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # 동일 이미지  페이지 추가
        left = Inches(0)
        top = Inches(0)

        slide.shapes.add_picture(picture_path, left, top)


# 파일 저장
prs.save(os.path.join(path_dir, "test_presentation.pptx"))

# resize and attatch
img_path = os.path.join(img_folder, img_files[0])
img = Image.open(img_path)
new_size = (1100, 825)  # 예시로 변경
resized_img = img.resize(new_size, Image.Resampling.LANCZOS)
print(resized_img.size)  # (1100, 825
# Crop the image before resizing
# Example: Crop the central part of the image
left = 250
upper = 100
right = left + 720
lower = upper + 590
crop_area = (left, upper, right, lower)
cropped_img = resized_img.crop(crop_area)
# print(cropped_img.size)  # (110, 155)
# 변경된 이미지 저장
cropped_img.save(os.path.join(path_dir, temp_img_path))
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 이미지 추가
left = Inches(0)
top = Inches(0)
slide.shapes.add_picture(str(os.path.join(path_dir, temp_img_path)), left, top)

prs.save(os.path.join(path_dir, "test_presentation.pptx"))

for img_name in img_files:
    img_path = os.path.join(img_folder, img_name)
    if os.path.isfile(img_path):
        # 이미지 열기
        with Image.open(img_path) as img:
            # 원본 이미지 비율 유지하면서 A4 크기에 맞게 조정
            new_size = (1100, 825)  # 예시로 800x600 픽셀로 변경
            resized_img = img.resize(new_size)

            # 변경된 이미지 저장
            # Crop the image before resizing
            # Example: Crop the central part of the image
            left = 250
            upper = 100
            right = left + 720
            lower = upper + 590
            crop_area = (left, upper, right, lower)
            cropped_img = resized_img.crop(crop_area)
            # print(cropped_img.size)  # (110, 155)
            # 변경된 이미지 저장
            cropped_img.save(os.path.join(path_dir, temp_img_path))

            # 조정된 이미지 저장 (임시 파일)

            # 새 슬라이드 추가 (빈 슬라이드)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 이미지 추가
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(os.path.join(path_dir, temp_img_path), left, top)

            os.remove(os.path.join(path_dir, temp_img_path))


# 파일 저장
prs.save(os.path.join(path_dir, "test_presentation.pptx"))
